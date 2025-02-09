from flask import Flask, request, jsonify
import os
from dotenv import load_dotenv
from flask_cors import CORS
from GetNews import GetNewsContent
from openai import OpenAI, OpenAIError
import base64
from pptx import Presentation
from PIL import Image
import io
import traceback
import PyPDF2
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import json

app = Flask(__name__)

# Simple CORS configuration
CORS(app, supports_credentials=True)

# Load environment variables
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
NEWS_API_KEY = os.getenv("NEWS_API_KEY")

UPLOAD_FOLDER = "uploads"
TEMP_DIR = "temp"  # Add this line

ALLOWED_EXTENSIONS = {"pdf", "pptx"}


if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
if not os.path.exists(TEMP_DIR): 
    os.makedirs(TEMP_DIR)  


def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


# Validate API keys
if not OPENAI_API_KEY:
    raise ValueError("OPENAI_API_KEY not found in environment variables")
if not NEWS_API_KEY:
    raise ValueError("NEWS_API_KEY not found in environment variables")

# Initialize OpenAI client with proper authentication
client = OpenAI(
    api_key=OPENAI_API_KEY,
)


def convert_pptx_to_image(pptx_path, slide_index=0):
    """Convert a PowerPoint slide to an image"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        prs = Presentation(pptx_path)
        if slide_index >= len(prs.slides):
            raise ValueError("Slide index out of range")

        slide = prs.slides[slide_index]
        
        # Get slide dimensions
        width = int(prs.slide_width * 96 / 914400)  # convert EMU to pixels
        height = int(prs.slide_height * 96 / 914400)
        
        # Create a new image with white background
        slide_image = Image.new('RGB', (width, height), 'white')
        draw = ImageDraw.Draw(slide_image)

        # Process all shapes on the slide
        for shape in slide.shapes:
            # Handle text boxes
            if hasattr(shape, "text"):
                left = int(shape.left * 96 / 914400)
                top = int(shape.top * 96 / 914400)
                try:
                    # Draw text (basic implementation)
                    draw.text((left, top), shape.text, fill='black')
                except Exception as e:
                    print(f"Error drawing text: {e}")

            # Handle pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = Image.open(io.BytesIO(shape.image.blob))
                    left = int(shape.left * 96 / 914400)
                    top = int(shape.top * 96 / 914400)
                    width = int(shape.width * 96 / 914400)
                    height = int(shape.height * 96 / 914400)
                    image = image.resize((width, height))
                    slide_image.paste(image, (left, top))
                except Exception as e:
                    print(f"Error processing image: {e}")

        # Save to buffer
        image_stream = io.BytesIO()
        slide_image.save(image_stream, format='PNG')
        image_stream.seek(0)
        return image_stream.getvalue()

    except Exception as e:
        print(f"Error converting PowerPoint to image: {str(e)}")
        traceback.print_exc()
        raise

def convert_pdf_to_image(pdf_path, page_number=0):
    """Convert a PDF page to an image"""
    print('here2')
    print(pdf_path)
    # Replace PdfFileReader with PdfReader
    pdf = PyPDF2.PdfReader(pdf_path)
    if page_number >= len(pdf.pages):
        raise ValueError("Page number out of range")

    # Create a temporary buffer to save the page as PNG
    image_stream = io.BytesIO()

    Image.new("RGB", (800, 600), "white").save(image_stream, format="PNG")
    image_stream.seek(0)
    return image_stream.getvalue()


def generate_lesson(file_path, prompt):
    try:
        # Convert PowerPoint slide to image
        print('here', file_path.rsplit(".", 1)[1].lower())

        if file_path.rsplit(".", 1)[1].lower() == "pptx":
            image_data = convert_pptx_to_image(file_path)
        elif file_path.rsplit(".", 1)[1].lower() == "pdf":
            image_data = convert_pdf_to_image(file_path)
        else:   
            raise Exception("Invalid file format")
        
        # Convert image data to base64
        base64_image = base64.b64encode(image_data).decode("utf-8")

        # Create message with file attachment
        response = client.chat.completions.create(
            model="gpt-4o-mini",  # Updated model name
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text", 
                            "text": prompt
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{base64_image}"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1000
        )

        # Extract and validate the response
        response_content = response.choices[0].message.content
        try:
            # Attempt to parse as JSON first
            json_response = json.loads(response_content)
            return json_response
        except json.JSONDecodeError:
            # If not valid JSON, create a structured response
            return [{
                "slide": 1,
                "suggestions": [{
                    "content": response_content,
                    "link": ""
                }]
            }]

    except FileNotFoundError:
        raise Exception("File not found")
    except OpenAIError as e:
        raise Exception(f"OpenAI API error: {str(e)}")
    except Exception as e:
        print(f"Error in generate_lesson: {str(e)}")
        traceback.print_exc()
        raise Exception(f"Error processing file: {str(e)}")


def generate_w_pdfs(file_path, prompt):
    try:
        # Read the PDF file
        with open(file_path, "rb") as file:
            # Create PDF reader object
            pdf_reader = PyPDF2.PdfReader(file)

            # Extract text from all pages
            text_content = ""
            for page in pdf_reader.pages:
                text_content += page.extract_text()

        # Create message with the extracted text
        response = client.chat.completions.create(
            model="gpt-4o-mini",  # Make sure to use an appropriate model
            messages=[
                {
                    "role": "user",
                    "content": f"{prompt}\n\nHere is the syllabus content:\n{text_content}",
                }
            ],
            max_tokens=1000,
        )

        return response.choices[0].message.content

    except FileNotFoundError:
        raise Exception("File not found")
    except OpenAIError as e:
        raise Exception(f"OpenAI API Error: {str(e)}")
    except Exception as e:
        raise Exception(f"Error processing file: {str(e)}")


@app.route("/lesson-plan", methods=["POST"])
def suggest():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded"}), 400      

        file = request.files["file"]
        subject = request.form.get("subject", "")

        if not file or not subject:
            return jsonify({"error": "File and subject are required"}), 400
        
        file_path = os.path.join(TEMP_DIR, file.filename)
        file.save(file_path)
        
        news = GetNewsContent(NEWS_API_KEY)
        subject_news = news.get_subject_news(subject, days_back=7)

        if not os.path.exists(file_path):
            return jsonify({"error": "File not found"}), 404

        prompt = "Analyze this lesson plan give me suggested changes based on these recent news articles: {subject_news}. Focus on incorporating current research trends and modern teaching methodologies in education. The changes will be returned in this format: { \"slide\": <slide_number>, \"suggestions\": [ { \"content\": <suggestion_text>, \"link\": <source_link> } ] }Only return json format"
        


        response = generate_lesson(file_path, prompt)
        # The response is already a Python object, no need to parse it again
        print(response[0])
        os.remove(file_path)
        return jsonify({"suggestion": [response[0]]})

    except OpenAIError as e:
        print(f"OpenAI API Error: {str(e)}")
        return jsonify({"error": "OpenAI API Error", "details": str(e)}), 401
    except Exception as e:
        print(f"Error in /suggest endpoint: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route("/convert-pptx", methods=["POST"])
def convert_pptx():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No selected file"}), 400

        # Save the file temporarily
        file_path = os.path.join(TEMP_DIR, file.filename)
        file.save(file_path)

        # Convert slides to images
        prs = Presentation(file_path)
        slides = []

        for i in range(len(prs.slides)):
            image_data = convert_pptx_to_image(file_path, i)
            base64_image = base64.b64encode(image_data).decode('utf-8')
            slides.append(f"data:image/png;base64,{base64_image}")

        # Clean up
        os.remove(file_path)

        return jsonify({"slides": slides})

    except Exception as e:
        print(f"Error converting PPTX: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route("/content-suggest", methods=["POST"])
def content_suggest():
    try:
        # Get the data from the request
        data = request.get_json()
        if not data or "filename" not in data or "subject" not in data:
            return jsonify({"error": "Filename and subject are required"}), 400

        # Get news for the specific subject
        news = GetNewsContent(NEWS_API_KEY)
        subject_news = news.get_subject_news(data["subject"], days_back=7)

        file_path = os.path.join(TEMP_DIR, data["filename"])  # Updated this line

        if not os.path.exists(file_path):
            return jsonify({"error": "File not found"}), 404

        prompt = (
            f"Analyze this syllabus and suggest improvements based on "
            f"these recent news articles in {data['subject']}: {subject_news}. "
            f"Focus on incorporating current research trends and "
            f"modern teaching methodologies in {data['subject']} education."
        )

        response = generate_w_pdfs(file_path, prompt)
        return jsonify({"suggestion": response})

    except OpenAIError as e:
        print(f"OpenAI API Error: {str(e)}")
        return jsonify({"error": "OpenAI API Error", "details": str(e)}), 401

    except Exception as e:
        print(f"Error in /content-suggest endpoint: {str(e)}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/upload", methods=["POST"])
def upload_file():
    try:
        print("Received request")
        print("Request headers:", dict(request.headers))
        print("Request files:", request.files)

        if "file" not in request.files:
            print("No file part in request")
            return {"error": "No file part"}, 400

        file = request.files["file"]
        subject = request.form.get("subject", "")

        if file.filename == "":
            print("No selected file")
            return {"error": "No selected file"}, 400

        if not subject:
            return {"error": "Subject is required"}, 400

        print(f"Processing file: {file.filename}")
        print(f"Subject: {subject}")

        # Save the file using TEMP_DIR constant
        file_path = os.path.join(TEMP_DIR, file.filename)
        file.save(file_path)
        print(f"File saved at: {file_path}")

        # Read and print the first few bytes of the file (for debugging)
        with open(file_path, "rb") as f:
            content = f.read(100)
        print(f"First 100 bytes: {content}")

        response = jsonify(
            {
                "message": "File processed successfully",
                "filename": file.filename,
                "subject": subject,
            }
        )
        return response

    except Exception as e:
        print(f"Error occurred: {str(e)}")
        traceback.print_exc()
        return {"error": str(e)}, 500


@app.route("/")
def home():
    return "Hello, Flask!"


@app.route("/about")
def about():
    return "About page"


@app.route("/mahin")
def mahin():
    return jsonify({"message": "Mahin page"})


if __name__ == "__main__":
    app.run(debug=True)
